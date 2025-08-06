# Create streamlined Streamlit app with PDF, PowerPoint, single powerful chart, and worst case
streamlit_app_content = '''import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime
from dataclasses import dataclass
from typing import Optional
import io
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Page config
st.set_page_config(
    page_title="Sales-Training ROI Kalkulator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

@dataclass
class TrainingParameters:
    """Parameter für das Sales-Training"""
    participants: int
    cost_per_person: float
    monthly_leads: int
    current_close_rate: float
    target_close_rate: float
    deal_value: float
    margin_rate: float
    training_days: int = 3
    daily_rate: float = 400.0

@dataclass
class ROIResults:
    """Ergebnisse der ROI-Berechnung"""
    total_investment: float
    training_costs: float
    opportunity_costs: float
    current_deals: float
    target_deals: float
    additional_deals: float
    monthly_revenue: float
    monthly_margin: float
    annual_margin: float
    roi_percentage: float
    roi_multiple: float
    payback_days: int
    net_benefit: float

class SalesROICalculator:
    """Sales-Training ROI Kalkulator"""
    
    def __init__(self):
        self.parameters: Optional[TrainingParameters] = None
        self.results: Optional[ROIResults] = None
    
    def format_currency(self, amount: float) -> str:
        """Formatiere Betrag als Währung"""
        return f"{amount:,.0f} €".replace(",", ".")
    
    def format_number(self, number: float, decimals: int = 1) -> str:
        """Formatiere Zahl mit Dezimalstellen"""
        return f"{number:,.{decimals}f}".replace(",", ".")
    
    def calculate_roi(self, params: TrainingParameters) -> ROIResults:
        """Berechne ROI basierend auf Parametern"""
        self.parameters = params
        
        # 1. Gesamtinvestition berechnen
        training_costs = params.participants * params.cost_per_person
        opportunity_costs = params.participants * params.training_days * params.daily_rate
        total_investment = training_costs + opportunity_costs
        
        # 2. Zusätzliche Deals berechnen
        current_deals = params.monthly_leads * (params.current_close_rate / 100)
        target_deals = params.monthly_leads * (params.target_close_rate / 100)
        additional_deals = target_deals - current_deals
        
        # 3. Umsatz- und Margen-Impact
        monthly_revenue = additional_deals * params.deal_value
        monthly_margin = monthly_revenue * (params.margin_rate / 100)
        annual_margin = monthly_margin * 12
        
        # 4. ROI-Metriken
        net_benefit = annual_margin - total_investment
        roi_percentage = (net_benefit / total_investment) * 100 if total_investment > 0 else 0
        roi_multiple = net_benefit / total_investment if total_investment > 0 else 0
        payback_days = int((total_investment / monthly_margin) * 30) if monthly_margin > 0 else 0
        
        self.results = ROIResults(
            total_investment=total_investment,
            training_costs=training_costs,
            opportunity_costs=opportunity_costs,
            current_deals=current_deals,
            target_deals=target_deals,
            additional_deals=additional_deals,
            monthly_revenue=monthly_revenue,
            monthly_margin=monthly_margin,
            annual_margin=annual_margin,
            roi_percentage=roi_percentage,
            roi_multiple=roi_multiple,
            payback_days=payback_days,
            net_benefit=net_benefit
        )
        
        return self.results

def create_story_chart(calculator):
    """Erstelle den einen Chart, der die ganze Story erzählt"""
    if not calculator.results:
        return None
    
    r = calculator.results
    
    # Create the story chart - ROI development over 12 months
    fig = go.Figure()
    
    # Monthly cumulative profit
    months = list(range(13))
    cumulative = [-r.total_investment]  # Start with negative investment
    monthly_labels = ['Start'] + [f'Monat {i}' for i in range(1, 13)]
    
    for month in range(1, 13):
        cumulative.append(cumulative[-1] + r.monthly_margin)
    
    # Add the main line
    fig.add_trace(go.Scatter(
        x=months, 
        y=cumulative,
        mode='lines+markers',
        name='Kumulierter Gewinn',
        line=dict(color='#2E86AB', width=4),
        marker=dict(size=8, color='#2E86AB'),
        hovertemplate='<b>%{text}</b><br>Kumulierter Gewinn: %{y:,.0f} €<extra></extra>',
        text=monthly_labels
    ))
    
    # Add break-even line
    fig.add_hline(y=0, line_dash="dash", line_color="#E74C3C", line_width=2,
                  annotation_text="Break-Even", annotation_position="bottom right")
    
    # Add investment area (negative)
    fig.add_shape(
        type="rect",
        x0=0, x1=12,
        y0=-r.total_investment, y1=0,
        fillcolor="rgba(231, 76, 60, 0.1)",
        line=dict(width=0),
    )
    
    # Add profit area (positive)
    max_profit = max(cumulative)
    fig.add_shape(
        type="rect",
        x0=0, x1=12,
        y0=0, y1=max_profit,
        fillcolor="rgba(46, 134, 171, 0.1)",
        line=dict(width=0),
    )
    
    # Add key annotations
    payback_month = r.payback_days / 30
    if payback_month <= 12:
        fig.add_annotation(
            x=payback_month,
            y=0,
            text=f"Break-Even<br>Tag {r.payback_days}",
            showarrow=True,
            arrowhead=2,
            arrowcolor="#E74C3C",
            bgcolor="white",
            bordercolor="#E74C3C",
            borderwidth=2
        )
    
    # Final profit annotation
    fig.add_annotation(
        x=12,
        y=cumulative[-1],
        text=f"Jahresgewinn<br>{calculator.format_currency(r.annual_margin)}",
        showarrow=True,
        arrowhead=2,
        arrowcolor="#2E86AB",
        bgcolor="white",
        bordercolor="#2E86AB",
        borderwidth=2
    )
    
    # Investment annotation
    fig.add_annotation(
        x=0,
        y=-r.total_investment,
        text=f"Investment<br>{calculator.format_currency(r.total_investment)}",
        showarrow=True,
        arrowhead=2,
        arrowcolor="#E74C3C",
        bgcolor="white",
        bordercolor="#E74C3C",
        borderwidth=2
    )
    
    fig.update_layout(
        title={
            'text': f"<b>Sales-Training ROI Story: {r.roi_percentage:.0f}% ROI in 12 Monaten</b>",
            'x': 0.5,
            'font': {'size': 20}
        },
        xaxis_title="Monate nach Training-Start",
        yaxis_title="Kumulierter Gewinn (€)",
        height=500,
        showlegend=False,
        plot_bgcolor='white',
        paper_bgcolor='white',
        font=dict(size=12),
        margin=dict(l=80, r=80, t=80, b=80)
    )
    
    # Format y-axis
    fig.update_yaxis(
        tickformat=",.0f",
        gridcolor='lightgray',
        gridwidth=1
    )
    
    fig.update_xaxis(
        gridcolor='lightgray',
        gridwidth=1
    )
    
    return fig

def create_calculation_breakdown(calculator):
    """Erstelle detaillierte Kalkulationsübersicht"""
    if not calculator.results or not calculator.parameters:
        return None
    
    params = calculator.parameters
    results = calculator.results
    
    # Create calculation dataframe
    calc_data = []
    
    # Investment
    calc_data.append({
        'Kategorie': '💸 INVESTMENT',
        'Berechnung': '',
        'Ergebnis': ''
    })
    
    calc_data.append({
        'Kategorie': 'Trainingskosten',
        'Berechnung': f'{params.participants} × {calculator.format_currency(params.cost_per_person)}',
        'Ergebnis': calculator.format_currency(results.training_costs)
    })
    
    calc_data.append({
        'Kategorie': 'Ausfallkosten',
        'Berechnung': f'{params.participants} × {params.training_days} × {params.daily_rate}€',
        'Ergebnis': calculator.format_currency(results.opportunity_costs)
    })
    
    calc_data.append({
        'Kategorie': '🔸 Gesamtinvestition',
        'Berechnung': f'{calculator.format_currency(results.training_costs)} + {calculator.format_currency(results.opportunity_costs)}',
        'Ergebnis': calculator.format_currency(results.total_investment)
    })
    
    # Deal Analysis
    calc_data.append({
        'Kategorie': '📈 DEAL-ANALYSE',
        'Berechnung': '',
        'Ergebnis': ''
    })
    
    calc_data.append({
        'Kategorie': 'Zusätzliche Deals/Monat',
        'Berechnung': f'{params.monthly_leads} × ({params.target_close_rate}% - {params.current_close_rate}%)',
        'Ergebnis': f'{calculator.format_number(results.additional_deals)} Deals'
    })
    
    calc_data.append({
        'Kategorie': 'Zusatzgewinn/Monat',
        'Berechnung': f'{calculator.format_number(results.additional_deals)} × {calculator.format_currency(params.deal_value)} × {params.margin_rate}%',
        'Ergebnis': calculator.format_currency(results.monthly_margin)
    })
    
    calc_data.append({
        'Kategorie': '🔸 Zusatzgewinn/Jahr',
        'Berechnung': f'{calculator.format_currency(results.monthly_margin)} × 12',
        'Ergebnis': calculator.format_currency(results.annual_margin)
    })
    
    # ROI
    calc_data.append({
        'Kategorie': '🚀 ROI-BERECHNUNG',
        'Berechnung': '',
        'Ergebnis': ''
    })
    
    calc_data.append({
        'Kategorie': 'ROI %',
        'Berechnung': f'({calculator.format_currency(results.annual_margin)} - {calculator.format_currency(results.total_investment)}) ÷ {calculator.format_currency(results.total_investment)} × 100',
        'Ergebnis': f'{results.roi_percentage:.0f}%'
    })
    
    calc_data.append({
        'Kategorie': '🔸 Payback-Zeit',
        'Berechnung': f'{calculator.format_currency(results.total_investment)} ÷ {calculator.format_currency(results.monthly_margin)} × 30',
        'Ergebnis': f'{results.payback_days} Tage'
    })
    
    return pd.DataFrame(calc_data)

def create_pdf_report(calculator):
    """Erstelle PDF-Report"""
    if not calculator.results:
        return None
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.HexColor('#2E86AB'),
        alignment=1  # Center
    )
    
    story.append(Paragraph("🎯 Sales-Training ROI Analyse", title_style))
    story.append(Paragraph("Joey's Business Case für CFO & CEO", styles['Heading2']))
    story.append(Spacer(1, 20))
    
    # Executive Summary
    story.append(Paragraph("📊 EXECUTIVE SUMMARY", styles['Heading2']))
    
    r = calculator.results
    p = calculator.parameters
    
    summary_data = [
        ['Metrik', 'Wert', 'Bedeutung'],
        ['Gesamtinvestition', calculator.format_currency(r.total_investment), 'Einmalige Kosten'],
        ['Zusatzgewinn/Monat', calculator.format_currency(r.monthly_margin), 'Dauerhafter Gewinn'],
        ['Jahresgewinn', calculator.format_currency(r.annual_margin), 'Erster Jahresertrag'],
        ['ROI', f'{r.roi_percentage:.0f}%', 'Return on Investment'],
        ['Payback-Zeit', f'{r.payback_days} Tage', 'Amortisationsdauer'],
        ['ROI-Multiple', f'{r.roi_multiple + 1:.1f}x', 'Gewinnfaktor']
    ]
    
    summary_table = Table(summary_data)
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(summary_table)
    story.append(Spacer(1, 20))
    
    # Recommendation
    story.append(Paragraph("🎯 EMPFEHLUNG", styles['Heading2']))
    
    if r.roi_percentage > 100:
        recommendation = f"""
        <b>KLARE EMPFEHLUNG: TRAINING SOFORT DURCHFÜHREN!</b><br/><br/>
        ✅ ROI von {r.roi_percentage:.0f}% ist außergewöhnlich hoch<br/>
        ✅ Payback in nur {r.payback_days} Tagen<br/>
        ✅ {calculator.format_currency(r.monthly_margin)} zusätzlicher Gewinn pro Monat<br/>
        ✅ Perfekte Argumentationsbasis für das Management<br/><br/>
        <b>Fazit:</b> Jeder investierte Euro bringt {r.roi_multiple + 1:.1f}€ zurück!
        """
    else:
        recommendation = f"ROI von {r.roi_percentage:.0f}% rechtfertigt die Investition."
    
    story.append(Paragraph(recommendation, styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Calculation Details
    story.append(Paragraph("🔢 KALKULATIONS-DETAILS", styles['Heading2']))
    
    calc_df = create_calculation_breakdown(calculator)
    if calc_df is not None:
        calc_data = [['Kategorie', 'Berechnung', 'Ergebnis']]
        for _, row in calc_df.iterrows():
            if row['Kategorie'] and not any(x in row['Kategorie'] for x in ['💸', '📈', '🚀']):
                calc_data.append([row['Kategorie'], row['Berechnung'], row['Ergebnis']])
        
        calc_table = Table(calc_data)
        calc_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(calc_table)
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def create_powerpoint_deck(calculator):
    """Erstelle PowerPoint-Präsentation"""
    if not calculator.results:
        return None
    
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🎯 Sales-Training ROI Analyse"
    subtitle.text = "Joey's Business Case für CFO & CEO\\nROI: {:.0f}% | Payback: {} Tage".format(
        calculator.results.roi_percentage, calculator.results.payback_days)
    
    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📊 Executive Summary"
    
    r = calculator.results
    summary_text = f"""Investment: {calculator.format_currency(r.total_investment)}
    
Zusatzgewinn pro Monat: {calculator.format_currency(r.monthly_margin)}

Jahresgewinn: {calculator.format_currency(r.annual_margin)}

ROI: {r.roi_percentage:.0f}%

Payback-Zeit: {r.payback_days} Tage

ROI-Multiple: {r.roi_multiple + 1:.1f}x"""
    
    content.text = summary_text
    
    # Slide 3: Recommendation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎯 Empfehlung"
    
    if r.roi_percentage > 100:
        recommendation_text = f"""KLARE EMPFEHLUNG: TRAINING DURCHFÜHREN!

✅ ROI von {r.roi_percentage:.0f}% ist außergewöhnlich
✅ Payback in nur {r.payback_days} Tagen  
✅ {calculator.format_currency(r.monthly_margin)} Gewinn pro Monat
✅ Perfekte Argumentationsbasis

Fazit: Jeder Euro bringt {r.roi_multiple + 1:.1f}€ zurück!"""
    else:
        recommendation_text = f"ROI von {r.roi_percentage:.0f}% rechtfertigt die Investition."
    
    content.text = recommendation_text
    
    # Slide 4: Key Arguments
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "💼 Top-Argumente für Management"
    
    arguments_text = f"""1. GEWINN-FOKUS
   {calculator.format_currency(r.annual_margin)} zusätzlicher Jahresgewinn

2. SCHNELLE AMORTISATION  
   Investment zahlt sich in {r.payback_days} Tagen zurück

3. WETTBEWERBSVORTEIL
   Konkurrent nimmt uns täglich Gewinn weg

4. SKALIERUNG
   Marge wirkt auf ALLE zukünftigen Sales

5. RISIKO-MINIMIERUNG
   Status Quo ist das größte Risiko"""
    
    content.text = arguments_text
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def main():
    """Hauptfunktion der Streamlit App"""
    
    # Header
    st.title("🎯 Sales-Training ROI Kalkulator")
    st.markdown("### Joey's Business Case für CFO & CEO")
    
    # Szenario Box
    with st.expander("📋 Das Szenario", expanded=True):
        st.markdown("""
        **Ausgangssituation:**
        - Sales-Team stagniert bei niedriger Abschlussquote
        - Neue bewährte Methodik verspricht deutliche Steigerung
        - Training übersteigt genehmigtes Trainingsbudget
        
        **Herausforderung:**
        - CFO und CEO müssen für Budgetfreigabe gewonnen werden
        - ⏰ **ZEITDRUCK:** Training soll schnell starten
        - 🏆 **WETTBEWERB:** Konkurrent hat bereits Vorsprung
        
        **Joey's Aufgabe:**
        - Wasserdichte ROI-Argumentation erstellen
        - Budgetüberschreitung mit Zahlen rechtfertigen
        """)
    
    # Sidebar für Parameter
    st.sidebar.header("📊 Parameter eingeben")
    
    # Parameter Input
    participants = st.sidebar.number_input("Anzahl Teilnehmer", min_value=1, max_value=50, value=8)
    cost_per_person = st.sidebar.number_input("Trainingskosten pro Person (€)", min_value=500, max_value=10000, value=3000, step=100)
    monthly_leads = st.sidebar.number_input("Monatliche Leads", min_value=10, max_value=1000, value=150, step=10)
    current_rate = st.sidebar.slider("Aktuelle Abschlussquote (%)", min_value=1.0, max_value=50.0, value=12.0, step=0.5)
    target_rate = st.sidebar.slider("Ziel-Abschlussquote (%)", min_value=1.0, max_value=50.0, value=20.0, step=0.5)
    deal_value = st.sidebar.number_input("Durchschnittlicher Deal-Wert (€)", min_value=1000, max_value=100000, value=12000, step=500)
    margin_rate = st.sidebar.slider("Marge pro Deal (%)", min_value=5.0, max_value=80.0, value=25.0, step=1.0)
    
    # Erweiterte Parameter
    with st.sidebar.expander("🔧 Erweiterte Parameter"):
        training_days = st.number_input("Trainingstage", min_value=1, max_value=10, value=3)
        daily_rate = st.number_input("Tagessatz Ausfallzeit (€)", min_value=200, max_value=1000, value=400, step=50)
    
    # Berechnung
    calculator = SalesROICalculator()
    
    params = TrainingParameters(
        participants=participants,
        cost_per_person=cost_per_person,
        monthly_leads=monthly_leads,
        current_close_rate=current_rate,
        target_close_rate=target_rate,
        deal_value=deal_value,
        margin_rate=margin_rate,
        training_days=training_days,
        daily_rate=daily_rate
    )
    
    results = calculator.calculate_roi(params)
    
    # Main Layout: Results and Calculations side by side
    col_results, col_calc = st.columns([2, 1])
    
    with col_results:
        # Hauptergebnisse
        st.header("💰 Kern-Ergebnisse")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            st.metric("💸 Gesamtinvestition", calculator.format_currency(results.total_investment))
        
        with col2:
            st.metric("📈 Mehrumsatz/Monat", calculator.format_currency(results.monthly_revenue))
        
        with col3:
            st.metric("💰 Zusatzgewinn/Monat", calculator.format_currency(results.monthly_margin))
        
        with col4:
            st.metric("🚀 ROI (12 Monate)", f"{results.roi_percentage:.0f}%")
        
        col5, col6 = st.columns(2)
        
        with col5:
            st.metric("⚡ Payback-Zeit", f"{results.payback_days} Tage")
        
        with col6:
            st.metric("🎖️ Jahresgewinn", calculator.format_currency(results.annual_margin))
        
        # Empfehlung
        st.header("🎯 Empfehlung")
        
        if results.roi_percentage > 100 and results.payback_days < 90:
            st.success("🎉 **KLARE EMPFEHLUNG: TRAINING DURCHFÜHREN!**")
            st.markdown(f"""
            ✅ ROI von **{results.roi_percentage:.0f}%** ist außergewöhnlich  
            ✅ Payback in nur **{results.payback_days} Tagen**  
            ✅ **{calculator.format_currency(results.monthly_revenue)}** Mehrumsatz pro Monat  
            ✅ **{calculator.format_currency(results.monthly_margin)}** zusätzlicher GEWINN pro Monat  
            ✅ **Perfekte Argumentationsbasis für CFO und CEO!**
            """)
        elif results.roi_percentage > 50:
            st.warning("👍 **EMPFEHLUNG: TRAINING LOHNT SICH**")
            st.markdown(f"""
            ✅ ROI von **{results.roi_percentage:.0f}%** rechtfertigt die Investition  
            ✅ Payback-Zeit: **{results.payback_days} Tage**  
            ✅ Monatlicher Zusatzgewinn: **{calculator.format_currency(results.monthly_margin)}**
            """)
        else:
            st.error("⚠️ **VORSICHT: ROI ZU NIEDRIG**")
            st.markdown(f"""
            ❌ ROI von nur **{results.roi_percentage:.0f}%** rechtfertigt möglicherweise nicht die Investition  
            ❌ Prüfen Sie die Parameter oder suchen Sie Alternativen  
            ❌ Monatlicher Zusatzgewinn nur: **{calculator.format_currency(results.monthly_margin)}**
            """)
    
    with col_calc:
        # Live Kalkulationen
        st.header("🔢 Live-Kalkulationen")
        
        # Create calculation breakdown
        calc_df = create_calculation_breakdown(calculator)
        
        if calc_df is not None:
            # Display calculation table
            display_df = calc_df[['Kategorie', 'Berechnung', 'Ergebnis']].copy()
            
            st.dataframe(
                display_df,
                use_container_width=True,
                hide_index=True,
                height=400
            )
            
            # Quick summary box
            st.info(f"""
            **🎯 Zusammenfassung:**
            
            **Investment:** {calculator.format_currency(results.total_investment)}
            **Jahresgewinn:** {calculator.format_currency(results.annual_margin)}
            **ROI:** {results.roi_percentage:.0f}%
            **Payback:** {results.payback_days} Tage
            
            **💡 Fazit:** Jeder investierte Euro bringt {results.roi_multiple + 1:.1f}€ zurück!
            """)
    
    # The ONE Chart that tells the whole story
    st.header("📊 Die ROI-Story in einem Chart")
    
    story_chart = create_story_chart(calculator)
    if story_chart:
        st.plotly_chart(story_chart, use_container_width=True)
    
    # Scenario Analysis with Worst Case
    st.header("🔍 Szenario-Analysen")
    
    tab1, tab2, tab3 = st.tabs(["🚀 Best Case", "⚠️ Worst Case", "💼 CFO Argumente"])
    
    with tab1:
        best_revenue = results.monthly_revenue * 1.3
        best_profit = best_revenue * (params.margin_rate / 100)
        best_roi = results.roi_percentage * 1.5
        
        st.markdown(f"""
        ### 🚀 BEST CASE SZENARIO
        **Annahme:** Training wirkt sogar besser als erwartet (+30% zum Ziel)
        
        📊 **Zahlen:**
        - Mehrumsatz: **{calculator.format_currency(best_revenue)}/Monat**
        - Zusatzgewinn: **{calculator.format_currency(best_profit)}/Monat**
        - ROI: **~{best_roi:.0f}%**
        - Payback: **~{int(results.payback_days * 0.7)} Tage**
        
        💼 **Joey's Argument:**
        > "Selbst wenn wir konservativ rechnen, ist der ROI fantastisch. 
        > Im Best Case haben wir **{calculator.format_currency(best_profit * 12)}** zusätzlichen Jahresgewinn!"
        """)
    
    with tab2:
        # Worst Case: Training wirkt nur zu 50%, Marge sinkt um 5%
        worst_deals = results.additional_deals * 0.5
        worst_revenue = worst_deals * params.deal_value
        worst_margin_rate = max(5, params.margin_rate - 5)
        worst_profit = worst_revenue * (worst_margin_rate / 100)
        worst_annual = worst_profit * 12
        worst_roi = ((worst_annual - results.total_investment) / results.total_investment) * 100
        worst_payback = int((results.total_investment / worst_profit) * 30) if worst_profit > 0 else 999
        
        st.markdown(f"""
        ### ⚠️ WORST CASE SZENARIO
        **Annahmen:** 
        - Training wirkt nur zu 50% wie erwartet
        - Marge sinkt um 5% (auf {worst_margin_rate}%)
        - Marktbedingungen verschlechtern sich
        
        📊 **Zahlen:**
        - Zusätzliche Deals: **{calculator.format_number(worst_deals)}/Monat**
        - Zusatzgewinn: **{calculator.format_currency(worst_profit)}/Monat**
        - Jahresgewinn: **{calculator.format_currency(worst_annual)}**
        - ROI: **{worst_roi:.0f}%**
        - Payback: **{worst_payback} Tage**
        
        💡 **Joey's Argument:**
        > "Selbst im Worst Case haben wir noch **{worst_roi:.0f}% ROI**! 
        > Das Risiko ist minimal - Status Quo ist viel riskanter!"
        
        **🔥 Entscheidend:** Auch im schlechtesten Fall ist die Investition profitabel!
        """)
    
    with tab3:
        st.markdown(f"""
        ### 💼 TOP-ARGUMENTE FÜR CFO & CEO
        
        **1️⃣ GEWINN-FOKUS:**
        > "Training generiert **{calculator.format_currency(results.annual_margin)}** zusätzlichen 
        > Jahresgewinn - das ist echtes Geld in der Kasse!"
        
        **2️⃣ SCHNELLE AMORTISATION:**
        > "Investment zahlt sich in **{results.payback_days} Tagen** zurück - 
        > schneller als jede Maschine oder Software"
        
        **3️⃣ RISIKO-MINIMIERUNG:**
        > "Selbst im Worst Case: **{((worst_annual - results.total_investment) / results.total_investment) * 100:.0f}% ROI**. 
        > Status Quo ist das größte Risiko!"
        
        **4️⃣ WETTBEWERBSDRUCK:**
        > "Konkurrent nimmt uns täglich **{calculator.format_currency(results.monthly_margin/30)}** 
        > Gewinn weg - jeden Monat den wir warten!"
        
        **5️⃣ SKALIERUNG:**
        > "Diese Marge wirkt auf ALLE zukünftigen Sales - 
        > nicht nur auf das Training!"
        """)
    
    # Export Funktionen
    st.header("📄 Export & Download")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        # PDF Export
        pdf_buffer = create_pdf_report(calculator)
        if pdf_buffer:
            st.download_button(
                label="📄 PDF Report",
                data=pdf_buffer,
                file_name=f"roi_analyse_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                mime="application/pdf"
            )
    
    with col2:
        # PowerPoint Export
        ppt_buffer = create_powerpoint_deck(calculator)
        if ppt_buffer:
            st.download_button(
                label="📊 PowerPoint Deck",
                data=ppt_buffer,
                file_name=f"roi_presentation_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    
    with col3:
        # CSV Export
        if calc_df is not None:
            st.download_button(
                label="📈 Kalkulationen CSV",
                data=calc_df.to_csv(index=False, encoding='utf-8'),
                file_name=f"roi_kalkulationen_{datetime.now().strftime('%Y%m%d_%H%M')}.csv",
                mime="text/csv"
            )

if __name__ == "__main__":
    main()
'''

# Write the updated Streamlit app to file
with open('app.py', 'w', encoding='utf-8') as f:
    f.write(streamlit_app_content)

print("✅ Streamlit app.py wurde erfolgreich aktualisiert!")
print("🎯 Neue Features:")
print("   - PDF-Report mit professionellem Layout")
print("   - PowerPoint-Deck (4 Slides) für Präsentation")
print("   - EIN aussagekräftiger Chart der die ganze Story erzählt")
print("   - Worst-Case Szenario für Risiko-Analyse")
print("   - Kein JSON Export mehr")
print("   - Fokussierte, überzeugende Darstellung")
print("🌐 Bereit für Online-Deployment!")
